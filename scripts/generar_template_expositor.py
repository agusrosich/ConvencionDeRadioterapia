"""
Genera el Template Expositor.pptx para la RT International Institute Convention 2026.

Crea un PPTX de 5 slides estandarizado para los expositores invitados,
manteniendo la identidad visual del Formato Base.pptx original.
"""

import os
import copy
import zipfile
import shutil
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from lxml import etree

# Paths
BASE_DIR = Path(__file__).resolve().parent.parent
FORMATO_BASE = BASE_DIR / "Presentaciones" / "Base de presentaciones" / "Formato Base.pptx"
LOGO_PATH = BASE_DIR / "Logos" / "CONVENCION 2026 (1).png"  # Gold on transparent
OUTPUT_PATH = BASE_DIR / "Presentaciones" / "Base de presentaciones" / "Template Expositor.pptx"

# Design constants (from original PPTX analysis)
BG_COLOR = RGBColor(0xEB, 0xEB, 0xE6)       # Beige background
TEXT_DARK = RGBColor(0x2B, 0x2B, 0x2B)       # Main text
BLUE_INST = RGBColor(0x30, 0x3A, 0x7F)       # Institutional blue
GOLD = RGBColor(0xB6, 0x90, 0x5A)            # Gold accent
GOLD_DARK = RGBColor(0xA3, 0x7B, 0x40)       # Dark gold

FONT_TITLE = "Arial Unicode MS"
FONT_BODY = "Mukti"

# Positions (EMU) - from original analysis
HEADER_LEFT_POS = (0, 672465)
HEADER_LEFT_SIZE = (5032023, 356235)
HEADER_RIGHT_POS = (13646280, 672465)
HEADER_RIGHT_SIZE = (3613020, 356235)
FOOTER_LEFT_POS = (1028700, 9061132)
FOOTER_LEFT_SIZE = (3613020, 356235)
LOGO_POS = (15452790, 8388832)
LOGO_SIZE = (1749448, 1749448)
LINE_Y = 7868354

# Content area
TITLE_POS = (1434135, 2097424)
TITLE_SIZE = (15400000, 900000)
CONTENT_POS = (1434135, 3100000)
CONTENT_SIZE = (15400000, 4500000)
SUBTITLE_POS = (1434135, 2900000)
SUBTITLE_SIZE = (15400000, 500000)


def add_background(slide):
    """Set solid background color."""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = BG_COLOR


def add_textbox(slide, left, top, width, height, text, font_name, font_size,
                font_color, bold=False, alignment=PP_ALIGN.LEFT, alpha=None):
    """Add a styled text box to a slide."""
    txBox = slide.shapes.add_textbox(Emu(left), Emu(top), Emu(width), Emu(height))
    tf = txBox.text_frame
    tf.word_wrap = True

    # Remove default padding
    tf.margin_top = 0
    tf.margin_bottom = 0
    tf.margin_left = 0
    tf.margin_right = 0

    p = tf.paragraphs[0]
    p.alignment = alignment

    # Set line spacing
    pPr = p._pPr
    if pPr is None:
        pPr = p._p.get_or_add_pPr()
    nsmap = {"a": "http://schemas.openxmlformats.org/drawingml/2006/main"}
    lnSpc = etree.SubElement(pPr, f"{{{nsmap['a']}}}lnSpc")
    spcPts = etree.SubElement(lnSpc, f"{{{nsmap['a']}}}spcPts")
    spcPts.set("val", str(int(font_size.pt * 140)))  # ~140% line spacing

    run = p.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = font_size
    run.font.bold = bold
    run.font.color.rgb = font_color

    if alpha is not None:
        # Add alpha to the color
        rPr = run._r.get_or_add_rPr()
        solidFill = rPr.find(f"{{{nsmap['a']}}}solidFill")
        if solidFill is not None:
            clr = solidFill.find(f"{{{nsmap['a']}}}srgbClr")
            if clr is not None:
                alpha_elem = etree.SubElement(clr, f"{{{nsmap['a']}}}alpha")
                alpha_elem.set("val", str(int(alpha * 1000)))

    return txBox


def add_header_footer(slide):
    """Add the standard header and footer elements to a slide."""
    # Header left: "CONVENTION 2026"
    add_textbox(slide, *HEADER_LEFT_POS, *HEADER_LEFT_SIZE,
                "CONVENTION 2026", FONT_BODY, Pt(21), GOLD,
                alignment=PP_ALIGN.CENTER, alpha=53)

    # Header right: "MARZO 2026"
    add_textbox(slide, *HEADER_RIGHT_POS, *HEADER_RIGHT_SIZE,
                "MARZO 2026", FONT_BODY, Pt(21), GOLD,
                alignment=PP_ALIGN.RIGHT, alpha=50)

    # Footer left: "@RTINTERNATIONAL"
    add_textbox(slide, *FOOTER_LEFT_POS, *FOOTER_LEFT_SIZE,
                "@RTINTERNATIONAL", FONT_BODY, Pt(21), GOLD_DARK)

    # Logo bottom-right
    if LOGO_PATH.exists():
        slide.shapes.add_picture(
            str(LOGO_PATH),
            Emu(LOGO_POS[0]), Emu(LOGO_POS[1]),
            Emu(LOGO_SIZE[0]), Emu(LOGO_SIZE[1])
        )

    # Horizontal line
    line = slide.shapes.add_shape(
        1,  # MSO_SHAPE.RECTANGLE
        Emu(-410910), Emu(LINE_Y),
        Emu(18698910), Emu(12700)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = GOLD
    line.line.fill.background()


def add_content_placeholder(slide, top=None, height=None, hint_text="Escriba su contenido aquí"):
    """Add content area with placeholder text."""
    t = top or CONTENT_POS[1]
    h = height or CONTENT_SIZE[1]
    add_textbox(slide, CONTENT_POS[0], t, CONTENT_SIZE[0], h,
                hint_text, FONT_BODY, Pt(20), RGBColor(0x99, 0x99, 0x99))


def create_slide_portada(prs):
    """Slide 1: Portada del expositor."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    add_background(slide)

    # Large title: specialty name
    add_textbox(slide, 1028700, 3200000, 16230600, 1600000,
                "ESPECIALIDAD", FONT_TITLE, Pt(100), TEXT_DARK,
                bold=True, alignment=PP_ALIGN.CENTER)

    # Case info
    add_textbox(slide, 4368979, 5200000, 9550042, 500000,
                "TÍTULO DE LA PRESENTACIÓN", FONT_BODY, Pt(24), TEXT_DARK,
                alignment=PP_ALIGN.CENTER)

    # Separator line (decorative)
    sep = slide.shapes.add_shape(
        1, Emu(7000000), Emu(5900000), Emu(4288000), Emu(25400)
    )
    sep.fill.solid()
    sep.fill.fore_color.rgb = GOLD
    sep.line.fill.background()

    # Presenter name
    add_textbox(slide, 4368979, 6100000, 9550042, 540000,
                "NOMBRE DEL EXPOSITOR/A", FONT_BODY, Pt(32), TEXT_DARK,
                alignment=PP_ALIGN.CENTER)

    # Institution
    add_textbox(slide, 4368979, 6700000, 9550042, 400000,
                "Institución / Especialidad", FONT_BODY, Pt(24),
                RGBColor(0x66, 0x66, 0x66), alignment=PP_ALIGN.CENTER)

    add_header_footer(slide)
    return slide


def create_slide_aporte(prs):
    """Slide 2: Aporte de la especialidad."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide)

    add_textbox(slide, *TITLE_POS, *TITLE_SIZE,
                "APORTE DE LA ESPECIALIDAD", FONT_TITLE, Pt(48),
                TEXT_DARK, bold=True)

    add_textbox(slide, SUBTITLE_POS[0], 2850000, *SUBTITLE_SIZE,
                "¿Cuál es la contribución específica de su disciplina en este caso?",
                FONT_BODY, Pt(18), RGBColor(0x80, 0x80, 0x80))

    add_content_placeholder(slide, top=3500000, height=4000000,
                           hint_text="• Describa el aporte de su especialidad\n• Hallazgos relevantes\n• Implicancias para el manejo del caso")

    add_header_footer(slide)
    return slide


def create_slide_decision(prs):
    """Slide 3: Punto de decisión clínica."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide)

    add_textbox(slide, *TITLE_POS, *TITLE_SIZE,
                "PUNTO DE DECISIÓN CLÍNICA", FONT_TITLE, Pt(48),
                TEXT_DARK, bold=True)

    add_textbox(slide, SUBTITLE_POS[0], 2850000, *SUBTITLE_SIZE,
                "¿Cuál es el punto de decisión clínica relevante?",
                FONT_BODY, Pt(18), RGBColor(0x80, 0x80, 0x80))

    add_content_placeholder(slide, top=3500000, height=4000000,
                           hint_text="• Identifique el punto de decisión clave\n• Opciones terapéuticas consideradas\n• Aspectos controversiales")

    add_header_footer(slide)
    return slide


def create_slide_conducta(prs):
    """Slide 4: Conducta propuesta."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide)

    add_textbox(slide, *TITLE_POS, *TITLE_SIZE,
                "CONDUCTA PROPUESTA", FONT_TITLE, Pt(48),
                TEXT_DARK, bold=True)

    add_textbox(slide, SUBTITLE_POS[0], 2850000, *SUBTITLE_SIZE,
                "¿Cuál es la conducta que propone y por qué?",
                FONT_BODY, Pt(18), RGBColor(0x80, 0x80, 0x80))

    add_content_placeholder(slide, top=3500000, height=4000000,
                           hint_text="• Conducta terapéutica propuesta\n• Justificación de la decisión\n• Resultados esperados")

    add_header_footer(slide)
    return slide


def create_slide_fundamento(prs):
    """Slide 5: Fundamento."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide)

    add_textbox(slide, *TITLE_POS, *TITLE_SIZE,
                "FUNDAMENTO", FONT_TITLE, Pt(48),
                TEXT_DARK, bold=True)

    add_textbox(slide, SUBTITLE_POS[0], 2850000, *SUBTITLE_SIZE,
                "Evidencia o experiencia institucional que sustenta la conducta",
                FONT_BODY, Pt(18), RGBColor(0x80, 0x80, 0x80))

    add_content_placeholder(slide, top=3500000, height=4000000,
                           hint_text="• Referencias clave (máximo 2-3)\n• Experiencia institucional relevante\n• Datos que soporten la decisión")

    add_header_footer(slide)
    return slide


def main():
    print("Generando Template Expositor.pptx...")

    # Create presentation with same dimensions as original
    prs = Presentation()
    prs.slide_width = Emu(18288000)   # 20 inches
    prs.slide_height = Emu(10287000)  # 11.25 inches

    # Create the 5 slides
    create_slide_portada(prs)
    create_slide_aporte(prs)
    create_slide_decision(prs)
    create_slide_conducta(prs)
    create_slide_fundamento(prs)

    # Save
    OUTPUT_PATH.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUTPUT_PATH))
    print(f"Template guardado en: {OUTPUT_PATH}")

    # Now embed the Mukti font from the original PPTX
    embed_fonts_from_original(str(FORMATO_BASE), str(OUTPUT_PATH))
    print("Fuentes embebidas copiadas del original.")
    print("¡Listo!")


def embed_fonts_from_original(source_pptx, target_pptx):
    """Copy embedded font data from source PPTX to target PPTX."""
    import tempfile

    temp_path = target_pptx + ".tmp"
    shutil.copy2(target_pptx, temp_path)

    with zipfile.ZipFile(source_pptx, 'r') as src_zip:
        font_files = [f for f in src_zip.namelist() if f.startswith('ppt/fonts/')]
        if not font_files:
            os.remove(temp_path)
            return

        with zipfile.ZipFile(temp_path, 'r') as tmp_zip:
            with zipfile.ZipFile(target_pptx, 'w', zipfile.ZIP_DEFLATED) as out_zip:
                # Update Content_Types to include font data
                content_types = tmp_zip.read('[Content_Types].xml').decode('utf-8')
                if 'fntdata' not in content_types:
                    content_types = content_types.replace(
                        '</Types>',
                        '<Default Extension="fntdata" ContentType="application/x-fontdata"/></Types>'
                    )

                # Copy all existing files from target (with updated Content_Types)
                for item in tmp_zip.namelist():
                    if item == '[Content_Types].xml':
                        out_zip.writestr(item, content_types)
                    else:
                        out_zip.writestr(item, tmp_zip.read(item))

                # Copy font files from source
                for font_file in font_files:
                    out_zip.writestr(font_file, src_zip.read(font_file))

    os.remove(temp_path)


if __name__ == "__main__":
    main()
